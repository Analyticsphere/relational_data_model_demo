from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

NAVY=RGBColor(0x16,0x32,0x4F); TEAL=RGBColor(0x2E,0x8B,0x8B); DARK=RGBColor(0x22,0x2A,0x33)
GREY=RGBColor(0x5A,0x63,0x6E); LIGHT=RGBColor(0xED,0xF1,0xF5); WHITE=RGBColor(0xFF,0xFF,0xFF)
ACCENT=RGBColor(0xCC,0x7D,0x15)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def notes(s,t): s.notes_slide.notes_text_frame.text=t
def titlebar(s,title,kicker=None):
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(1.15))
    bar.fill.solid();bar.fill.fore_color.rgb=NAVY;bar.line.fill.background();bar.shadow.inherit=False
    tf=bar.text_frame;tf.margin_left=Inches(0.5);tf.margin_top=Inches(0.12);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0];r=p.add_run();r.text=title;r.font.size=Pt(28);r.font.bold=True;r.font.color.rgb=WHITE;r.font.name="Calibri"
    if kicker:
        p2=tf.add_paragraph();rr=p2.add_run();rr.text=kicker;rr.font.size=Pt(13);rr.font.color.rgb=RGBColor(0xBB,0xD5,0xD7);rr.font.italic=True
    acc=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(1.15),SW,Inches(0.06))
    acc.fill.solid();acc.fill.fore_color.rgb=TEAL;acc.line.fill.background();acc.shadow.inherit=False
def textbox(s,text,left,top,width,height,size=18,color=DARK,bold=False,align=PP_ALIGN.LEFT,italic=False):
    tb=s.shapes.add_textbox(left,top,width,height);tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.alignment=align;r=p.add_run();r.text=text
    r.font.size=Pt(size);r.font.color.rgb=color;r.font.bold=bold;r.font.italic=italic;r.font.name="Calibri";return tb
def bullets(s,items,left,top,width,height,size=18,gap=10):
    tb=s.shapes.add_textbox(left,top,width,height);tf=tb.text_frame;tf.word_wrap=True
    norm=[(x,) if isinstance(x,str) else x for x in items]
    for i,(txt,*lvl) in enumerate(norm):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(gap);p.level=(lvl[0] if lvl else 0)
        r=p.add_run();r.text=("•  "+txt) if p.level==0 else ("–  "+txt)
        r.font.size=Pt(size if p.level==0 else size-2);r.font.color.rgb=DARK if p.level==0 else GREY;r.font.name="Calibri"
    return tb
def table(s,data,left,top,width,height,col_widths=None,fs=13,hdr_fs=13):
    gt=s.shapes.add_table(len(data),len(data[0]),left,top,width,height);t=gt.table
    if col_widths:
        for c,w in enumerate(col_widths):t.columns[c].width=w
    for ri,row in enumerate(data):
        for ci,val in enumerate(row):
            cell=t.cell(ri,ci)
            cell.margin_left=Inches(0.08);cell.margin_right=Inches(0.08);cell.margin_top=Inches(0.03);cell.margin_bottom=Inches(0.03)
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE;tf=cell.text_frame;tf.word_wrap=True
            p=tf.paragraphs[0];r=p.add_run();r.text=str(val);r.font.name="Calibri"
            if ri==0:
                cell.fill.solid();cell.fill.fore_color.rgb=NAVY;r.font.color.rgb=WHITE;r.font.bold=True;r.font.size=Pt(hdr_fs)
            else:
                cell.fill.solid();cell.fill.fore_color.rgb=WHITE if ri%2 else LIGHT
                r.font.color.rgb=DARK;r.font.size=Pt(fs)
                if ci==0:r.font.bold=True
    return t
def picture(s,path,top,max_w,max_h,left=None):
    iw,ih=Image.open(path).size;ar=iw/ih;w=max_w;h=Emu(int(w/ar))
    if h>max_h:h=max_h;w=Emu(int(h*ar))
    if left is None:left=Emu(int((SW-w)/2))
    s.shapes.add_picture(path,left,top,width=w,height=h)

CODEBG=RGBColor(0x1E,0x29,0x35); CODEFG=RGBColor(0xE6,0xED,0xF3); CODECM=RGBColor(0x8B,0xA6,0xC0)
def codebox(s,code,left,top,width,height,size=12):
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,left,top,width,height)
    box.fill.solid(); box.fill.fore_color.rgb=CODEBG; box.line.color.rgb=TEAL; box.shadow.inherit=False
    tf=box.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.2); tf.margin_right=Inches(0.15); tf.margin_top=Inches(0.12); tf.vertical_anchor=MSO_ANCHOR.TOP
    for i,line in enumerate(code.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.0
        r=p.add_run(); r.text=line or " "; r.font.name="Courier New"; r.font.size=Pt(size)
        r.font.color.rgb=CODECM if line.strip().startswith(("--","#")) else CODEFG
    return box

A="docs/archive/slide_assets/"

# 1 TITLE
s=slide()
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH);bg.fill.solid();bg.fill.fore_color.rgb=NAVY;bg.line.fill.background();bg.shadow.inherit=False
band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(4.35),SW,Inches(0.08));band.fill.solid();band.fill.fore_color.rgb=TEAL;band.line.fill.background();band.shadow.inherit=False
textbox(s,"Connect Relational Data Model",Inches(0.9),Inches(2.3),Inches(11.5),Inches(1.4),size=44,color=WHITE,bold=True)
textbox(s,"Why a data model — the accepted Dictionary-Direct model + its enhancement roadmap for PR2",Inches(0.9),Inches(3.5),Inches(11.5),Inches(0.8),size=20,color=RGBColor(0xBB,0xD5,0xD7))
textbox(s,"Team walk-through  ·  detailed notes in internal_pitch.md",Inches(0.9),Inches(4.6),Inches(11.5),Inches(0.5),size=14,color=RGBColor(0x9A,0x9D,0xA3),italic=True)
notes(s,"We already believe in this idea (OMOP) — we just never applied it to our own Connect data. The accepted model is Dictionary-Direct; further capabilities are adopted as incremental enhancements. See internal_pitch.md — The ask.")

# 1b WHY A DATA MODEL — the plain-English value props (opener)
s=slide();titlebar(s,"Why a data model?","Make the data explain itself — so analysis doesn't depend on tribal knowledge or bespoke code")
CARDS=[
 ("Self-describing","Meaning travels with the data — no decoder ring required."),
 ("No tribal knowledge","Labels, structure & valid values live in the data — query it correctly without insider context."),
 ("Metadata in the data","Not in side-car spreadsheets or scripts that drift out of sync."),
 ("Standardized analyses","The same question is answered the same way by everyone."),
 ("A schema contract","Tools built on it keep running when new data arrives."),
 ("Shared abstractions","Every select-all handled one way; every loop one way — no per-question code."),
]
cw, gap, x0 = Inches(3.91), Inches(0.3), Inches(0.5)
rows_y = [Inches(2.05), Inches(4.15)]; ch = Inches(1.9)
for i,(t,d) in enumerate(CARDS):
    rr, cc = divmod(i, 3)
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(x0)+cc*(int(cw)+int(gap))), rows_y[rr], cw, ch)
    card.fill.solid(); card.fill.fore_color.rgb=LIGHT; card.line.color.rgb=TEAL; card.shadow.inherit=False
    tf=card.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.18); tf.margin_right=Inches(0.18); tf.margin_top=Inches(0.14); tf.vertical_anchor=MSO_ANCHOR.TOP
    p=tf.paragraphs[0]; r=p.add_run(); r.text=t; r.font.bold=True; r.font.size=Pt(16); r.font.color.rgb=NAVY; r.font.name="Calibri"
    p2=tf.add_paragraph(); p2.space_before=Pt(5); r2=p2.add_run(); r2.text=d; r2.font.size=Pt(11.5); r2.font.color.rgb=GREY; r2.font.name="Calibri"
band=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(6.35),Inches(12.33),Inches(0.72))
band.fill.solid(); band.fill.fore_color.rgb=NAVY; band.line.fill.background(); band.shadow.inherit=False
tf=band.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.margin_left=Inches(0.25)
p=tf.paragraphs[0]; r=p.add_run()
r.text="Today that knowledge lives in column-name conventions, spreadsheets, and analysts' heads. A data model puts it in the data."
r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name="Calibri"
notes(s,"The opener: why model at all. Six plain-English value props — self-describing; queryable without tribal knowledge; metadata in the data (not sidecar files); standardized analyses; a schema contract tools can build on (and survive new data); shared abstractions (all select-all handled one way, all loops one way). Sets up the next slide: this is exactly the OMOP lesson. See internal_pitch.md — Why a data model at all?")

# 2 WHY — the OMOP lesson (builds on the value props)
s=slide();titlebar(s,"…and it's the OMOP lesson — applied to our own data","The power isn't the table layout — it's relationships defined as data")
textbox(s,"Defined relationships between concepts & domains → analytics written once, reused across teams. That is exactly why we adopted OMOP.",Inches(0.5),Inches(1.4),Inches(12.3),Inches(0.7),size=17,color=TEAL,bold=True)
table(s,[["What OMOP defines","The same idea in our model"],
 ["concept — analyze standard concepts, not source codes","concept-ID spine — query concepts, not d_ columns"],
 ["concept_relationship / concept_ancestor — link & roll up","concept_relationship / question_equivalence (planned)"],
 ["long event tables + _source_concept_id","long responses fact + retained source columns"],
 ["standard model → tools reused across 200 sites","stable contract → shared views/tools across teams + PR2"]],
 Inches(0.6),Inches(2.25),Inches(12.1),Inches(3.0),col_widths=[Inches(5.9),Inches(6.2)],fs=14,hdr_fs=14)
textbox(s,"One source, not many — no cross-site network effect; reuse comes from the relationships, across our teams & researchers. OMOP becomes a downstream export.",Inches(0.6),Inches(6.4),Inches(12.1),Inches(0.7),size=12,color=GREY,italic=True)
notes(s,"Lead with the OMOP principle the team already accepts. See internal_pitch.md — Why a data model at all?")

# 3 PROBLEM
s=slide();titlebar(s,"The problem: we analyze raw source data","Our analysis-ready data is wide tables of opaque concept-ID columns")
bullets(s,["Dancing schema — every new answer/option/loop adds columns; pipelines break",
 "Not generically queryable — every analysis hardcodes column names",
 "Version drift — v1/v2 columns reconciled by hand",
 "Ambiguous missingness — blank = not selected? not shown? not taken?",
 "No built-in governance — PHI/PII is a manual, unenforced allow-list"],
 Inches(0.7),Inches(1.6),Inches(11.8),Inches(3.6),size=20,gap=14)
textbox(s,"In OMOP terms, these wide tables are source data — and we're analyzing them raw.",Inches(0.7),Inches(6.2),Inches(11.8),Inches(0.7),size=16,color=ACCENT,bold=True)
notes(s,"The five pains. Punchline lands after the exhibits. See internal_pitch.md — Why — the problems we all live with.")

# 4 EXHIBIT 1
s=slide();titlebar(s,"Exhibit 1 — our own Module 1 report","\"Merged Module 1 Summary Statistics\" · 6,234 lines · scheduled pipeline")
textbox(s,"~650 lines rebuild the model by hand before the first statistic:",Inches(0.7),Inches(1.45),Inches(11.8),Inches(0.6),size=18,bold=True,color=TEAL)
bullets(s,["Hand-rolled v1/v2 merge — re-runs every refresh, column-name dependent",
 "150-entry label dictionary typed by hand — with duplicate-key bugs & typos",
 "Skip logic & loops reimplemented as bespoke functions",
 "Every missing value → \"Skipped this Question\" — 3 cases conflated (counts inflated)"],
 Inches(0.7),Inches(2.2),Inches(11.8),Inches(3.0),size=18,gap=12)
textbox(s,"→ The model does the merge once (GROUP BY concept), generates labels, and makes skip/loops data.",Inches(0.7),Inches(6.2),Inches(11.8),Inches(0.7),size=15,color=GREY,italic=True)
notes(s,"A colleague's flagship report — evidence, not criticism. See internal_pitch.md — This isn't hypothetical.")

# 5 EXHIBIT 2
s=slide();titlebar(s,"Exhibit 2 — the QA/QC engine","1,271-line engine + 14 Excel workbooks = 7,025 hand-written rules")
table(s,[["The rule checks…","# (all 14 workbooks)","Already in the model as"],
 ["Cross-variable condition","3,454  (49%)","skip_logic"],
 ["Numeric / date / time","883  (13%)","variable_metadata.variable_type"],
 ["String length","792  (11%)","variable_metadata.variable_length"],
 ["Answer ∈ valid option set","773  (11%)","response_options"],
 ["Populated / required","65  (1%)","skip_logic / is_required"],
 ["Genuinely custom","1,058  (15%)","— needs authoring"]],
 Inches(0.6),Inches(1.5),Inches(12.1),Inches(3.7),col_widths=[Inches(4.6),Inches(3.3),Inches(4.2)],fs=14,hdr_fs=14)
textbox(s,"Most re-state metadata the model centralizes; complex cross-variable logic is authored once — not re-coded in 3 places.",Inches(0.6),Inches(6.15),Inches(12.1),Inches(0.55),size=15,color=ACCENT,bold=True)
textbox(s,"Honest limits: ~35% are pure value/type/length checks; the skip-gated rest needs the richer skip-logic model (~8% complex tail authored once). Module 4 ≈ 50% custom.",Inches(0.6),Inches(6.72),Inches(12.1),Inches(0.55),size=11,color=GREY,italic=True)
notes(s,"Recalibrated: ~35% directly generatable (value/type/length); a further ~50% skip-gated -> becomes shared skip-logic data authored ONCE instead of re-coded across 3 repos (simple majority machine-generated, ~8% complex tail like percentDiff preserved as raw expression). ~15% genuinely custom. Contingent on the richer (QuickQ) skip_rule representation. Module 4 ~50% custom is the honest counter-case. See internal_pitch.md — QA/QC engine; CLAUDE.md — Skip-logic complexity & coverage.")

# 6 EXHIBIT 3
s=slide();titlebar(s,"Exhibit 3 — geocoding & unharmonized concept IDs","Same field, no relationship linking its variants")
bullets(s,["\"Street name of a residence\" = ~27 unrelated concept IDs (home ×11, seasonal ×10, work, school, childhood…)",
 "Three separate codebases each rebuild the same address crosswalk",
 "One even string-matches question labels as join keys (fragile — labels drift)",
 "No \"address\" or \"street_name\" concept generalizes → nothing reusable"],
 Inches(0.7),Inches(1.6),Inches(11.9),Inches(3.2),size=19,gap=13)
textbox(s,"Fix: concept_relationship / question_equivalence — author \"these are the same field\" once, as data.",Inches(0.7),Inches(6.1),Inches(11.9),Inches(0.8),size=16,color=ACCENT,bold=True)
notes(s,"Strongest case for the concept-equivalence enhancement — argues to pull it forward. See internal_pitch.md — geocoding.")

# 7 THE IDEA + full model ERD
s=slide();titlebar(s,"The idea — answers become rows","The model: the dictionary as-is + one long responses table")
picture(s,A+"model_a_full.png",Inches(1.45),Inches(12.7),Inches(5.05))
textbox(s,"New answers / options / loops are rows, never new columns. Full model ERD (Dictionary-Direct).",Inches(0.5),Inches(6.85),Inches(12.3),Inches(0.5),size=12,color=GREY,italic=True,align=PP_ALIGN.CENTER)
notes(s,"The model: CIDTool dictionary as-is + one responses fact. See internal_pitch.md — The model.")

# 7b A CLOSER LOOK: select-all (ease-in)
s=slide();titlebar(s,"A closer look: select-all questions","Same data, a clearer shape — the dictionary stays the source of truth")
textbox(s,"Today's encoding faithfully mirrors how the data is stored. The model just adds a logical view on top — nothing is lost, the dictionary is unchanged.",Inches(0.55),Inches(1.35),Inches(12.2),Inches(0.7),size=15,color=TEAL,bold=True)
table(s,[["Example: \"Have you lost any permanent teeth?\"","Dictionary today","Model — logical view"],
 ["The question","shifted into \"Source Question\"","one question (multi_select)"],
 ["Each option (e.g. \"from accident\")","promoted to its own \"Question\"","a response option"],
 ["The answer","a synthetic Yes/No per option","one row per option actually selected"],
 ["Add an option next year","a new column","a new row"]],
 Inches(0.6),Inches(2.25),Inches(12.1),Inches(3.0),col_widths=[Inches(4.7),Inches(3.7),Inches(3.7)],fs=13.5,hdr_fs=13)
textbox(s,"Phased & low-risk: the dictionary tables are untouched; this is a view we layer on — adopt it where it helps, keep everything else as-is.",Inches(0.6),Inches(6.35),Inches(12.1),Inches(0.7),size=15,color=GREY,italic=True)
notes(s,"EASE-IN framing for a skeptical audience. The current encoding isn't wrong — it faithfully mirrors physical storage (a multi-select really is N 0/1 columns). The model's improvement is logical: an option is an answer choice, not a question, so it belongs in response_options; 'selected' becomes the presence of a row instead of a synthetic Yes/No. Same tooth-loss example as the docs (899251483; options 812107266 accident / 452438775 decay / 886864375 other / 551489317 none; Yes/No = 353358909/104430631). Quest disambiguates select-all [ ] vs single-select ( ) vs grid |grid|. Nothing is thrown away; the dictionary remains the source of truth and this view sits on top. See CLAUDE.md — Source Question is overloaded.")

# 8 INCREMENTAL ENHANCEMENTS ROADMAP
s=slide();titlebar(s,"Incremental enhancements — the roadmap","Each a bounded add-on to the accepted model — same responses fact, dictionary stays the source of truth")
table(s,[["#","Enhancement","Value"],
 ["1","Normalized question-type view","templated per-type SQL across surveys"],
 ["2","Typed value columns","numeric analysis; coded answers join labels"],
 ["3","Improved version handling","unify V1/V2; offered-vs-not-selected answerable"],
 ["4","skip_logic (from Quest)","skip logic queryable; most QA/QC rules generated"],
 ["5","response_sessions","true missingness (not-asked vs not-answered)"],
 ["6","Concept equivalence (demo built)","harmonize reused fields once, as data"],
 ["7","Governance (sensitivity + IAM)","external-release readiness (PHI/PII) for PR2"],
 ["8","Analytics marts (dbt)","curated, reproducible derived variables + lineage"],
 ["9","Event plane","long-format win for operational/biospecimen data"]],
 Inches(0.6),Inches(1.45),Inches(12.1),Inches(4.7),col_widths=[Inches(0.55),Inches(4.55),Inches(7.0)],fs=12.5,hdr_fs=13)
textbox(s,"Adopt one at a time as need pulls it in. Governance (#7) gates external sharing; dbt marts (#8) deliver the most researcher value. See enhancement_backlog.md.",Inches(0.6),Inches(6.4),Inches(12.1),Inches(0.7),size=14,color=ACCENT,bold=True)
notes(s,"The larger redesign was not adopted wholesale; its features are these incremental enhancements on the accepted Dictionary-Direct model. Each attaches to the same responses fact and keeps the dictionary as the source of truth. Honest dependencies (from the backlog): CIDTool maturity for dims, a Quest parser for skip_logic, and governance as an org/IRB effort. See docs/enhancement_backlog.md.")

# 9b EVENTS GO LONG TOO (DevOps proposal + compatibility ask)
s=slide();titlebar(s,"New: events are going long-format too","DevOps is moving follow-ups from nested columns to per-type event tables")
bullets(s,["Per-type long tables: activities · collections · kits · surveys · incentives · refusals (+ collectionDetails)",
 "Shared key on every table: (Connect ID, Round) — Round + activity type = the encounter",
 "collectionDetails: ~120 nested fields → 14 columns; a new round/specimen = zero schema change",
 "Risk: human-readable names & values dropped the concept IDs → won't link to Primary/Secondary Source"],
 Inches(0.7),Inches(1.55),Inches(11.9),Inches(3.4),size=18,gap=13)
textbox(s,"Keep the long format — re-attach a concept_id to every column & value (one vocab, one ID: Mouthwash = MW = Home Mouthwash).",Inches(0.7),Inches(6.1),Inches(11.9),Inches(0.8),size=16,color=ACCENT,bold=True)
notes(s,"Same long-format thesis as our model, arrived at independently. The one fix: concept IDs, not strings. See CLAUDE.md — Event plane: DevOps long-format proposal; docs/devops_event_tables_memo.md.")

# 9c DRAFT EVENT PLANE diagram
s=slide();titlebar(s,"Draft: the reconciled event plane","Round = the encounter; surveys & events share it; concept-typed columns keep dictionary links")
picture(s,A+"event_plane.png",Inches(1.4),Inches(11.0),Inches(5.35))
textbox(s,"Separate per-type event facts + a unified view (the SMDB summary). DRAFT — pending DevOps confirmation of Round-as-encounter.",Inches(0.5),Inches(6.95),Inches(12.3),Inches(0.4),size=11,color=GREY,italic=True,align=PP_ALIGN.CENTER)
notes(s,"Round hub (~ OMOP visit) off participants; response_sessions = DevOps 'surveys' table; collection/kit/incentive/refusal facts key on (connect_id, round); v_participant_events = SMDB unified view. See dbml/data_model_events.dbml.")

# 9d CURATED DERIVED VARIABLES (dbt marts)
s=slide();titlebar(s,"Curated derived variables (dbt marts)","Define each research variable once — tested, governed, lineage to source")
textbox(s,"The payoff of \"write analytics once\": a canonical pack-years / BMI / risk score every researcher reuses — not re-derived per study.",Inches(0.55),Inches(1.35),Inches(12.2),Inches(0.7),size=15,color=TEAL,bold=True)
table(s,[["Group","Example marts","Canonical derived variable"],
 ["Behavioral","smoking · alcohol · physical_activity · diet","pack-years · MET-hours/week"],
 ["Anthropometric","anthropometry","BMI (height + weight)"],
 ["Host","reproductive_history · family_history · genetic_risk","parity · hereditary-risk flag · polygenic risk score"],
 ["Clinical","medical_history · medications · screening_history","comorbidity index · regular aspirin use · up-to-date-with-screening"],
 ["Environmental","environmental_exposures","Area Deprivation Index (from geocoded address)"]],
 Inches(0.6),Inches(2.25),Inches(12.1),Inches(2.8),col_widths=[Inches(2.3),Inches(5.6),Inches(4.2)],fs=12.5,hdr_fs=13)
textbox(s,"Why dbt: model-level lineage to source · tests = contracts in CI · one-directional layer boundary · sensitivity tiers flow to marts (enforced in IAM).",Inches(0.6),Inches(6.25),Inches(12.1),Inches(0.55),size=13,color=ACCENT,bold=True)
textbox(s,"Honest: derivation is real epi work (owner + sign-off per variable); marts are curated, not the only access — raw responses stay reachable.",Inches(0.6),Inches(6.8),Inches(12.1),Inches(0.5),size=11,color=GREY,italic=True)
notes(s,"This is the analytics-marts enhancement (#8), downstream of the model via dbt. Pitch = highest-value case of the OMOP 'write analytics once' thesis: derived variables defined once, canonically, with lineage — vs. re-coded per study (the 3 pain repos). dbt features map to our principles: dbt docs DAG = lineage to source (#12); SQL is the definition (no black box, #12); source()/ref()+access enforce Core→Analytic→Marts (#8); tests move much of the 7,025-rule QC into CI; sensitivity_tier flows via meta, enforced in IAM (#11); model versions/exposures = a derived-variable catalog. Recommend one mart PER construct (grouped), not a wide mega-table. Grain: participant, or participant×wave for longitudinal. Examples map to established cancer-prevention exposure domains (WCRF/AICR). See CLAUDE.md — Derived variables, lineage & dbt / Mart catalog.")

# 10 QUERIES
s=slide();titlebar(s,"Same queries, two ways","From painful/impossible → routine")
table(s,[["Standard query","Wide (today)","The model","With an enhancement"],
 ["Multi-select distribution","unpivot + v1/v2 COALESCE","filtered group-by; versions pool","plain group-by / view"],
 ["Labeled distribution, any question","bespoke CASE, no reuse","parameterized by concept_id","precomputed mart aggregate"],
 ["Completion & true missingness","ambiguous","still ambiguous","sessions + skip_logic (#5/#4)"],
 ["Non-PHI extract (governance)","manual allow-list","manual allow-list","sensitivity tier + IAM (#7)"],
 ["Harmonized field across surveys","rebuild crosswalk (3 repos)","pool a reused concept","concept equivalence (#6)"]],
 Inches(0.4),Inches(1.5),Inches(12.5),Inches(3.6),col_widths=[Inches(3.5),Inches(3.0),Inches(3.0),Inches(3.0)],fs=12.5,hdr_fs=13)
textbox(s,"The model fixes the loud, everyday pains; enhancements unlock missingness, governance & harmonization.",Inches(0.4),Inches(6.4),Inches(12.5),Inches(0.6),size=15,color=GREY,italic=True)
notes(s,"Worked SQL for each row is in internal_pitch.md — Value proposition (Q1–Q5).")

# 10b ANALYZING LONG FORMAT — demo 1: codes + labels side by side
s=slide();titlebar(s,"Analyzing long format (1/3) — labels are one join away","Every row says what it means: concept id + human label, side by side")
codebox(s,
"-- one join attaches a label to every answer\n"
"SELECT connect_id, question_text,\n"
"       response_concept_id, response_label, value\n"
"FROM responses\n"
"JOIN question USING (question_concept_id)\n"
"LEFT JOIN response USING (response_concept_id)\n"
"WHERE connect_id = 1001;",
Inches(0.5),Inches(1.45),Inches(6.1),Inches(2.7),size=12.5)
table(s,[["question_text","resp. id","response_label","value"],
 ["Sex","536341288","Female",""],
 ["Age","","","47"],
 ["Smoking Status","700000002","Former",""],
 ["Education","875342283","Bachelor's Degree",""],
 ["Have you lost…teeth?","812107266","Yes, from accident or injury",""],
 ["Have you lost…teeth?","452438775","Yes, from tooth decay or disease",""]],
 Inches(6.8),Inches(1.45),Inches(6.0),Inches(2.7),col_widths=[Inches(1.7),Inches(1.1),Inches(2.6),Inches(0.6)],fs=10.5,hdr_fs=10.5)
textbox(s,"No decoder ring, no 2,360-column map — the dictionary travels with the data.",Inches(0.5),Inches(6.4),Inches(12.3),Inches(0.5),size=14,color=ACCENT,bold=True)
notes(s,"Demo data is synthetic (sql/demo_long_format.sql, runs in DuckDB ~ BigQuery). The labels come from joining the dictionary; analysts never memorize concept ids. v_long packages this join so it's literally SELECT * FROM v_long.")

# 10c demo 2: pivot a subset back to wide — generic, in SQL / Python / R
s=slide();titlebar(s,"Analyzing long format (2/3) — wide is one generic PIVOT away","Pick your tool — one line, and none of them name a single column")
codebox(s,
"-- SQL (DuckDB) — column names come from the DATA, not from you\n"
"PIVOT (FROM v_long SELECT connect_id, question_text, answer\n"
"       WHERE question_type <> 'multi_select')\n"
"  ON question_text USING any_value(answer);\n"
"\n"
"# Python (pandas)\n"
"long.pivot(index='connect_id', columns='question_text', values='answer')\n"
"\n"
"# R (tidyr)\n"
"pivot_wider(long, id_cols=connect_id, names_from=question_text, values_from=answer)",
Inches(0.5),Inches(1.38),Inches(12.3),Inches(2.95),size=12)
table(s,[["connect_id","Age","Education","Sex","Smoking Status"],
 ["1001","47","Bachelor's Degree","Female","Former"],
 ["1002","62","High School Graduate or GED","Male","Current"],
 ["1003","55","Advanced Degree","Female","Never"],
 ["1004","39","Bachelor's Degree","Male","Never"],
 ["1005","71","High School Graduate or GED","Female","Former"]],
 Inches(1.6),Inches(4.55),Inches(10.1),Inches(1.7),col_widths=[Inches(1.5),Inches(1.0),Inches(4.3),Inches(1.5),Inches(1.8)],fs=11.5,hdr_fs=11.5)
textbox(s,"Column names come from the data, not from you — add a question, get a new column. Wrap it once as get_wide(survey) and every team reuses it.",Inches(0.5),Inches(6.4),Inches(12.3),Inches(0.5),size=13.5,color=ACCENT,bold=True)
textbox(s,"pandas & tidyr discover columns automatically; BigQuery's PIVOT lists the values (or generate them with dynamic SQL). `long` = the single-value answers.",Inches(0.5),Inches(6.92),Inches(12.3),Inches(0.4),size=10.5,color=GREY,italic=True)
notes(s,"Reframed: the pivot is GENERIC and metadata-driven — no hardcoded column list (the opposite of a wide-table script). NOTE on completeness: DuckDB's PIVOT implicitly groups by every column not in ON/USING, so you must project to (id, name, value) first — that's the subquery shown; this exact statement runs. The fully-dynamic ON…USING form is DuckDB; BigQuery's PIVOT needs the value list (FOR question_text IN (...)) or dynamic SQL. pandas.pivot and tidyr::pivot_wider are fully dynamic and are what analysts already use. All three verified against sql/demo_long_format.sql (synthetic). Select-all is excluded (multi-valued → belongs in long, next slide). Parameterize by survey → a reusable get_wide() helper (tool built on the contract).")

# 10d demo 3: counts are a simple GROUP BY (SATA is easier in long)
s=slide();titlebar(s,"Analyzing long format (3/3) — counts are a simple GROUP BY","And select-all gets *easier*: one group-by, no summing 0/1 columns")
codebox(s,
"-- how many selected each tooth-loss reason?\n"
"SELECT response_label, COUNT(*) AS n\n"
"FROM v_long\n"
"WHERE question_concept_id = 899251483\n"
"GROUP BY response_label\n"
"ORDER BY n DESC;",
Inches(0.5),Inches(1.5),Inches(6.0),Inches(2.4),size=12.5)
table(s,[["tooth_loss_reason","n"],
 ["Yes, from tooth decay or disease","2"],
 ["Yes, from accident or injury","2"],
 ["Yes, for some other reason","1"],
 ["No, I haven't lost any teeth","1"]],
 Inches(6.8),Inches(1.5),Inches(6.0),Inches(2.0),col_widths=[Inches(4.7),Inches(1.3)],fs=11.5,hdr_fs=11.5)
textbox(s,"In wide, that means summing N ever-changing 0/1 indicator columns; in long it's the same group-by you'd write for any question.",Inches(0.5),Inches(4.3),Inches(12.3),Inches(0.6),size=13,color=GREY,italic=True)
textbox(s,"Swap the concept id and the same one-liner counts smoking, education, anything — one pattern for every question.",Inches(0.5),Inches(6.45),Inches(12.3),Inches(0.5),size=14,color=ACCENT,bold=True)
notes(s,"Counts are GROUP BY — often simpler than wide. Select-all is the standout: wide makes you know and sum a changing set of indicator columns; long is one group-by on response_label. Same query shape works for every question by swapping the concept id (e.g. Smoking Status -> Never 2 / Former 2 / Current 1).")

# 11 RECOMMENDATION
s=slide();titlebar(s,"Recommendation")
bullets(s,["The Dictionary-Direct model is accepted — stand it up (low cost, low risk, immediately useful)",
 "Adopt enhancements incrementally — prioritize governance (#7, gates external sharing) and dbt analytics marts (#8, most researcher value)"],
 Inches(0.8),Inches(1.8),Inches(11.7),Inches(2.4),size=21,gap=18)
box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(4.6),Inches(11.7),Inches(1.4))
box.fill.solid();box.fill.fore_color.rgb=LIGHT;box.line.color.rgb=TEAL;box.shadow.inherit=False
tf=box.text_frame;tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE;tf.margin_left=Inches(0.3)
p=tf.paragraphs[0];r=p.add_run();r.text="The model is the accepted foundation.   Enhancements make it the governed, shareable product PR2 needs."
r.font.size=Pt(18);r.font.bold=True;r.font.color.rgb=NAVY;r.font.name="Calibri"
notes(s,"Ask: stand up the accepted model now; adopt enhancements incrementally (esp. governance #7 and dbt marts #8); consider pulling the concept-equivalence enhancement forward (geocoding). See internal_pitch.md — Recommendation.")

out="docs/archive/Connect_Data_Model_Pitch.pptx";prs.save(out)
print("saved",out,"with",len(prs.slides._sldIdLst),"slides")
