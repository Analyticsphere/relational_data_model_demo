from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

NAVY=RGBColor(0x16,0x32,0x4F); TEAL=RGBColor(0x2E,0x8B,0x8B); DARK=RGBColor(0x22,0x2A,0x33)
GREY=RGBColor(0x5A,0x63,0x6E); LIGHT=RGBColor(0xED,0xF1,0xF5); WHITE=RGBColor(0xFF,0xFF,0xFF)
ACCENT=RGBColor(0xCC,0x7D,0x15)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def notes(s,t): s.notes_slide.notes_text_frame.text=t
def titlebar(s,title,kicker=None):
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(1.15))
    bar.fill.solid();bar.fill.fore_color.rgb=NAVY;bar.line.fill.background();bar.shadow.inherit=False
    tf=bar.text_frame;tf.margin_left=Inches(0.5);tf.margin_top=Inches(0.12);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0];r=p.add_run();r.text=title;r.font.size=Pt(28);r.font.bold=True;r.font.color.rgb=WHITE;r.font.name="Calibri"
    if kicker:
        p2=tf.add_paragraph();rr=p2.add_run();rr.text=kicker;rr.font.size=Pt(13);rr.font.color.rgb=RGBColor(0xBB,0xD5,0xD7);rr.font.italic=True
    acc=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(1.15),SW,Inches(0.06))
    acc.fill.solid();acc.fill.fore_color.rgb=TEAL;acc.line.fill.background();acc.shadow.inherit=False
def textbox(s,text,left,top,width,height,size=18,color=DARK,bold=False,align=PP_ALIGN.LEFT,italic=False):
    tb=s.shapes.add_textbox(left,top,width,height);tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.alignment=align;r=p.add_run();r.text=text
    r.font.size=Pt(size);r.font.color.rgb=color;r.font.bold=bold;r.font.italic=italic;r.font.name="Calibri";return tb
def bullets(s,items,left,top,width,height,size=18,gap=10):
    tb=s.shapes.add_textbox(left,top,width,height);tf=tb.text_frame;tf.word_wrap=True
    norm=[(x,) if isinstance(x,str) else x for x in items]
    for i,(txt,*lvl) in enumerate(norm):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(gap);p.level=(lvl[0] if lvl else 0)
        r=p.add_run();r.text=("•  "+txt) if p.level==0 else ("–  "+txt)
        r.font.size=Pt(size if p.level==0 else size-2);r.font.color.rgb=DARK if p.level==0 else GREY;r.font.name="Calibri"
    return tb
def table(s,data,left,top,width,height,col_widths=None,fs=13,hdr_fs=13):
    gt=s.shapes.add_table(len(data),len(data[0]),left,top,width,height);t=gt.table
    if col_widths:
        for c,w in enumerate(col_widths):t.columns[c].width=w
    for ri,row in enumerate(data):
        for ci,val in enumerate(row):
            cell=t.cell(ri,ci)
            cell.margin_left=Inches(0.08);cell.margin_right=Inches(0.08);cell.margin_top=Inches(0.03);cell.margin_bottom=Inches(0.03)
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE;tf=cell.text_frame;tf.word_wrap=True
            p=tf.paragraphs[0];r=p.add_run();r.text=str(val);r.font.name="Calibri"
            if ri==0:
                cell.fill.solid();cell.fill.fore_color.rgb=NAVY;r.font.color.rgb=WHITE;r.font.bold=True;r.font.size=Pt(hdr_fs)
            else:
                cell.fill.solid();cell.fill.fore_color.rgb=WHITE if ri%2 else LIGHT
                r.font.color.rgb=DARK;r.font.size=Pt(fs)
                if ci==0:r.font.bold=True
    return t
def picture(s,path,top,max_w,max_h,left=None):
    iw,ih=Image.open(path).size;ar=iw/ih;w=max_w;h=Emu(int(w/ar))
    if h>max_h:h=max_h;w=Emu(int(h*ar))
    if left is None:left=Emu(int((SW-w)/2))
    s.shapes.add_picture(path,left,top,width=w,height=h)

A="docs/slide_assets/"

# 1 TITLE
s=slide()
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH);bg.fill.solid();bg.fill.fore_color.rgb=NAVY;bg.line.fill.background();bg.shadow.inherit=False
band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(4.35),SW,Inches(0.08));band.fill.solid();band.fill.fore_color.rgb=TEAL;band.line.fill.background();band.shadow.inherit=False
textbox(s,"Connect Relational Data Model",Inches(0.9),Inches(2.3),Inches(11.5),Inches(1.4),size=44,color=WHITE,bold=True)
textbox(s,"Why a data model — and a two-phase path to the PR2 research warehouse",Inches(0.9),Inches(3.5),Inches(11.5),Inches(0.8),size=20,color=RGBColor(0xBB,0xD5,0xD7))
textbox(s,"Team walk-through  ·  detailed notes in internal_pitch.md",Inches(0.9),Inches(4.6),Inches(11.5),Inches(0.5),size=14,color=RGBColor(0x9A,0x9D,0xA3),italic=True)
notes(s,"We already believe in this idea (OMOP) — we just never applied it to our own Connect data. Two phases: a fast Phase 1 win, then the governed PR2 warehouse. See internal_pitch.md — The ask.")

# 2 WHY
s=slide();titlebar(s,"Why a data model at all?","The power isn't the table layout — it's relationships defined as data")
textbox(s,"Defined relationships between concepts & domains → analytics written once, reused across teams. That is exactly why we adopted OMOP.",Inches(0.5),Inches(1.4),Inches(12.3),Inches(0.7),size=17,color=TEAL,bold=True)
table(s,[["What OMOP defines","The same idea in our model"],
 ["concept — analyze standard concepts, not source codes","concept-ID spine — query concepts, not d_ columns"],
 ["concept_relationship / concept_ancestor — link & roll up","concept_relationship / question_equivalence (planned)"],
 ["long event tables + _source_concept_id","long responses fact + retained source columns"],
 ["standard model → tools reused across 200 sites","stable contract → shared views/tools across teams + PR2"]],
 Inches(0.6),Inches(2.25),Inches(12.1),Inches(3.0),col_widths=[Inches(5.9),Inches(6.2)],fs=14,hdr_fs=14)
textbox(s,"One source, not many — no cross-site network effect; reuse comes from the relationships, across our teams & researchers. OMOP becomes a downstream export.",Inches(0.6),Inches(6.4),Inches(12.1),Inches(0.7),size=12,color=GREY,italic=True)
notes(s,"Lead with the OMOP principle the team already accepts. See internal_pitch.md — Why a data model at all?")

# 3 PROBLEM
s=slide();titlebar(s,"The problem: we analyze raw source data","Our analysis-ready data is wide tables of opaque concept-ID columns")
bullets(s,["Dancing schema — every new answer/option/loop adds columns; pipelines break",
 "Not generically queryable — every analysis hardcodes column names",
 "Version drift — v1/v2 columns reconciled by hand",
 "Ambiguous missingness — blank = not selected? not shown? not taken?",
 "No built-in governance — PHI/PII is a manual, unenforced allow-list"],
 Inches(0.7),Inches(1.6),Inches(11.8),Inches(3.6),size=20,gap=14)
textbox(s,"In OMOP terms, these wide tables are source data — and we're analyzing them raw.",Inches(0.7),Inches(6.2),Inches(11.8),Inches(0.7),size=16,color=ACCENT,bold=True)
notes(s,"The five pains. Punchline lands after the exhibits. See internal_pitch.md — Why — the problems we all live with.")

# 4 EXHIBIT 1
s=slide();titlebar(s,"Exhibit 1 — our own Module 1 report","\"Merged Module 1 Summary Statistics\" · 6,234 lines · scheduled pipeline")
textbox(s,"~650 lines rebuild the model by hand before the first statistic:",Inches(0.7),Inches(1.45),Inches(11.8),Inches(0.6),size=18,bold=True,color=TEAL)
bullets(s,["Hand-rolled v1/v2 merge — re-runs every refresh, column-name dependent",
 "150-entry label dictionary typed by hand — with duplicate-key bugs & typos",
 "Skip logic & loops reimplemented as bespoke functions",
 "Every missing value → \"Skipped this Question\" — 3 cases conflated (counts inflated)"],
 Inches(0.7),Inches(2.2),Inches(11.8),Inches(3.0),size=18,gap=12)
textbox(s,"→ The model does the merge once (GROUP BY concept), generates labels, and makes skip/loops data.",Inches(0.7),Inches(6.2),Inches(11.8),Inches(0.7),size=15,color=GREY,italic=True)
notes(s,"A colleague's flagship report — evidence, not criticism. See internal_pitch.md — This isn't hypothetical.")

# 5 EXHIBIT 2
s=slide();titlebar(s,"Exhibit 2 — the QA/QC engine","1,271-line engine + 14 Excel workbooks = 7,025 hand-written rules")
table(s,[["The rule checks…","# (all 14 workbooks)","Already in the model as"],
 ["Cross-variable condition","3,454  (49%)","skip_logic"],
 ["Numeric / date / time","883  (13%)","variable_metadata.variable_type"],
 ["String length","792  (11%)","variable_metadata.variable_length"],
 ["Answer ∈ valid option set","773  (11%)","response_options"],
 ["Populated / required","65  (1%)","skip_logic / is_required"],
 ["Genuinely custom","1,058  (15%)","— needs authoring"]],
 Inches(0.6),Inches(1.5),Inches(12.1),Inches(3.7),col_widths=[Inches(4.6),Inches(3.3),Inches(4.2)],fs=14,hdr_fs=14)
textbox(s,"Most re-state metadata the model centralizes; complex cross-variable logic is authored once — not re-coded in 3 places.",Inches(0.6),Inches(6.15),Inches(12.1),Inches(0.55),size=15,color=ACCENT,bold=True)
textbox(s,"Honest limits: ~35% are pure value/type/length checks; the skip-gated rest needs the richer skip-logic model (~8% complex tail authored once). Module 4 ≈ 50% custom.",Inches(0.6),Inches(6.72),Inches(12.1),Inches(0.55),size=11,color=GREY,italic=True)
notes(s,"Recalibrated: ~35% directly generatable (value/type/length); a further ~50% skip-gated -> becomes shared skip-logic data authored ONCE instead of re-coded across 3 repos (simple majority machine-generated, ~8% complex tail like percentDiff preserved as raw expression). ~15% genuinely custom. Contingent on the richer (QuickQ) skip_rule representation. Module 4 ~50% custom is the honest counter-case. See internal_pitch.md — QA/QC engine; CLAUDE.md — Skip-logic complexity & coverage.")

# 6 EXHIBIT 3
s=slide();titlebar(s,"Exhibit 3 — geocoding & unharmonized concept IDs","Same field, no relationship linking its variants")
bullets(s,["\"Street name of a residence\" = ~27 unrelated concept IDs (home ×11, seasonal ×10, work, school, childhood…)",
 "Three separate codebases each rebuild the same address crosswalk",
 "One even string-matches question labels as join keys (fragile — labels drift)",
 "No \"address\" or \"street_name\" concept generalizes → nothing reusable"],
 Inches(0.7),Inches(1.6),Inches(11.9),Inches(3.2),size=19,gap=13)
textbox(s,"Fix: concept_relationship / question_equivalence — author \"these are the same field\" once, as data.",Inches(0.7),Inches(6.1),Inches(11.9),Inches(0.8),size=16,color=ACCENT,bold=True)
notes(s,"Strongest case for the concept-equivalence plane; sits beyond Phase 2's first cut — argues to pull it forward. See internal_pitch.md — geocoding.")

# 7 THE IDEA + full Phase 1 ERD
s=slide();titlebar(s,"The idea — answers become rows","Phase 1: the dictionary as-is + one long responses table")
picture(s,A+"model_a_full.png",Inches(1.45),Inches(12.7),Inches(5.05))
textbox(s,"New answers / options / loops are rows, never new columns. Full Phase 1 ERD (Dictionary-Direct).",Inches(0.5),Inches(6.85),Inches(12.3),Inches(0.5),size=12,color=GREY,italic=True,align=PP_ALIGN.CENTER)
notes(s,"Phase 1 model: CIDTool dictionary as-is + one responses fact. Same thesis carries to Phase 2. See internal_pitch.md — The idea / Phase 1.")

# 7b A CLOSER LOOK: select-all (ease-in)
s=slide();titlebar(s,"A closer look: select-all questions","Same data, a clearer shape — the dictionary stays the source of truth")
textbox(s,"Today's encoding faithfully mirrors how the data is stored. The model just adds a logical view on top — nothing is lost, the dictionary is unchanged.",Inches(0.55),Inches(1.35),Inches(12.2),Inches(0.7),size=15,color=TEAL,bold=True)
table(s,[["Example: \"Have you lost any permanent teeth?\"","Dictionary today","Model — logical view"],
 ["The question","shifted into \"Source Question\"","one question (multi_select)"],
 ["Each option (e.g. \"from accident\")","promoted to its own \"Question\"","a response option"],
 ["The answer","a synthetic Yes/No per option","one row per option actually selected"],
 ["Add an option next year","a new column","a new row"]],
 Inches(0.6),Inches(2.25),Inches(12.1),Inches(3.0),col_widths=[Inches(4.7),Inches(3.7),Inches(3.7)],fs=13.5,hdr_fs=13)
textbox(s,"Phased & low-risk: the dictionary tables are untouched; this is a view we layer on — adopt it where it helps, keep everything else as-is.",Inches(0.6),Inches(6.35),Inches(12.1),Inches(0.7),size=15,color=GREY,italic=True)
notes(s,"EASE-IN framing for a skeptical audience. The current encoding isn't wrong — it faithfully mirrors physical storage (a multi-select really is N 0/1 columns). The model's improvement is logical: an option is an answer choice, not a question, so it belongs in response_options; 'selected' becomes the presence of a row instead of a synthetic Yes/No. Same tooth-loss example as the docs (899251483; options 812107266 accident / 452438775 decay / 886864375 other / 551489317 none; Yes/No = 353358909/104430631). Quest disambiguates select-all [ ] vs single-select ( ) vs grid |grid|. Nothing is thrown away; the dictionary remains the source of truth and this view sits on top. See CLAUDE.md — Source Question is overloaded.")

# 8 TWO PHASES table
s=slide();titlebar(s,"A two-phase path","Phase 1 proves it cheaply · Phase 2 builds the governed warehouse")
table(s,[["","Phase 1 — Dictionary-Direct","Phase 2 — PR2 warehouse"],
 ["What","dictionary as-is + 1 responses table","cleaned dims, placement bridge, sessions, governance"],
 ["Buys","stable schema, generic SQL, labels, v1/v2 pool","true missingness, governed access, marts, view library"],
 ["Lift","low — reuses CleanConnect + dictionary","larger, phased — builds on Phase 1's fact"],
 ["Release-ready","internal / analyst only","yes — external PR2 sharing"]],
 Inches(0.6),Inches(1.7),Inches(12.1),Inches(3.4),col_widths=[Inches(1.7),Inches(5.2),Inches(5.2)],fs=14,hdr_fs=14)
textbox(s,"Same responses fact both phases — Phase 1 is a down payment, not throwaway work.",Inches(0.6),Inches(5.6),Inches(12.1),Inches(0.6),size=15,color=GREY,italic=True)
notes(s,"Phase 2 adds the placement bridge, sessions, version-scoped options, layers, governance. See internal_pitch.md — Phase 1 / Phase 2.")

# 9 PHASE 2 full ERD
s=slide();titlebar(s,"Phase 2 — the full model","The researcher-facing, governed warehouse for PR2")
picture(s,A+"model_b_full.png",Inches(1.4),Inches(12.5),Inches(5.5))
textbox(s,"Placement bridge · sessions · version-scoped options · governance · Core→Analytic→Marts.",Inches(0.5),Inches(7.0),Inches(12.3),Inches(0.4),size=11,color=GREY,italic=True,align=PP_ALIGN.CENTER)
notes(s,"Full Phase 2 ERD. Dense by design — the SVG (docs/connect_data_model.svg) zooms cleanly. See internal_pitch.md — Phase 2.")

# 9a IS PHASE 2 DOABLE? (feasibility / risks)
s=slide();titlebar(s,"Is Phase 2 doable today?","Yes — technically achievable on our existing artifacts. The risks are dependencies, not the schema.")
textbox(s,"De-risked already: builds on Phase 1's same fact · CleanConnect did the cleaning · sessions derivable · hardest structures need no new tables.",Inches(0.55),Inches(1.35),Inches(12.2),Inches(0.7),size=14,color=TEAL,bold=True)
table(s,[["Honest risk","Why","How we de-risk"],
 ["CIDTool maturity","dimensions come from it; it's still an in-dev tool","audit its output now / be ready to build the transform"],
 ["Quest parser","skip logic, order, loops — no parser exists yet","prototype on module 1 first"],
 ["Governance","IRB / date-shift / suppression / IAM — org, not just code","start as a parallel workstream early"],
 ["Cross-team timing","couples to PR2 + CleanConnect + DevOps events","sequence around them; keep Phase 1 independent"]],
 Inches(0.6),Inches(2.25),Inches(12.1),Inches(2.95),col_widths=[Inches(2.7),Inches(5.0),Inches(4.4)],fs=12.5,hdr_fs=13)
textbox(s,"So Phase 2 is a funded, multi-quarter, multi-stakeholder program — not a single sprint. Which is exactly why we start with Phase 1.",Inches(0.6),Inches(6.4),Inches(12.1),Inches(0.6),size=15,color=ACCENT,bold=True)
notes(s,"Verdict: doable, nothing needs a capability Connect lacks — but it's dependencies + org/policy, not modeling. CIDTool verified 2026-06-22 as a small in-dev JS tool, not a production dimension emitter. Sequence: dims+placement+sessions first (mechanical); governance in parallel early (long pole, gates external release); concept plane/marts/OMOP deferred until pulled. BigQuery has no FK enforcement — relationships logical, enforced by transform + dbt tests. See CLAUDE.md — Phase 2 feasibility.")

# 9b EVENTS GO LONG TOO (DevOps proposal + compatibility ask)
s=slide();titlebar(s,"New: events are going long-format too","DevOps is moving follow-ups from nested columns to per-type event tables")
bullets(s,["Per-type long tables: activities · collections · kits · surveys · incentives · refusals (+ collectionDetails)",
 "Shared key on every table: (Connect ID, Round) — Round + activity type = the encounter",
 "collectionDetails: ~120 nested fields → 14 columns; a new round/specimen = zero schema change",
 "Risk: human-readable names & values dropped the concept IDs → won't link to Primary/Secondary Source"],
 Inches(0.7),Inches(1.55),Inches(11.9),Inches(3.4),size=18,gap=13)
textbox(s,"Keep the long format — re-attach a concept_id to every column & value (one vocab, one ID: Mouthwash = MW = Home Mouthwash).",Inches(0.7),Inches(6.1),Inches(11.9),Inches(0.8),size=16,color=ACCENT,bold=True)
notes(s,"Same long-format thesis as our model, arrived at independently. The one fix: concept IDs, not strings. See CLAUDE.md — Event plane: DevOps long-format proposal; docs/devops_event_tables_memo.md.")

# 9c DRAFT EVENT PLANE diagram
s=slide();titlebar(s,"Draft: the reconciled event plane","Round = the encounter; surveys & events share it; concept-typed columns keep dictionary links")
picture(s,A+"event_plane.png",Inches(1.4),Inches(11.0),Inches(5.35))
textbox(s,"Separate per-type event facts + a unified view (the SMDB summary). DRAFT — pending DevOps confirmation of Round-as-encounter.",Inches(0.5),Inches(6.95),Inches(12.3),Inches(0.4),size=11,color=GREY,italic=True,align=PP_ALIGN.CENTER)
notes(s,"Round hub (~ OMOP visit) off participants; response_sessions = DevOps 'surveys' table; collection/kit/incentive/refusal facts key on (connect_id, round); v_participant_events = SMDB unified view. See dbml/data_model_events.dbml.")

# 9d CURATED DERIVED VARIABLES (dbt marts)
s=slide();titlebar(s,"Curated derived variables (dbt marts)","Define each research variable once — tested, governed, lineage to source")
textbox(s,"The payoff of \"write analytics once\": a canonical pack-years / BMI / risk score every researcher reuses — not re-derived per study.",Inches(0.55),Inches(1.35),Inches(12.2),Inches(0.7),size=15,color=TEAL,bold=True)
table(s,[["Group","Example marts","Canonical derived variable"],
 ["Behavioral","smoking · alcohol · physical_activity · diet","pack-years · MET-hours/week"],
 ["Anthropometric","anthropometry","BMI (height + weight)"],
 ["Host","reproductive_history · family_history · genetic_risk","parity · hereditary-risk flag · polygenic risk score"],
 ["Clinical","medical_history · medications · screening_history","comorbidity index · regular aspirin use · up-to-date-with-screening"],
 ["Environmental","environmental_exposures","Area Deprivation Index (from geocoded address)"]],
 Inches(0.6),Inches(2.25),Inches(12.1),Inches(2.8),col_widths=[Inches(2.3),Inches(5.6),Inches(4.2)],fs=12.5,hdr_fs=13)
textbox(s,"Why dbt: model-level lineage to source · tests = contracts in CI · one-directional layer boundary · sensitivity tiers flow to marts (enforced in IAM).",Inches(0.6),Inches(6.25),Inches(12.1),Inches(0.55),size=13,color=ACCENT,bold=True)
textbox(s,"Honest: derivation is real epi work (owner + sign-off per variable); marts are curated, not the only access — raw responses stay reachable.",Inches(0.6),Inches(6.8),Inches(12.1),Inches(0.5),size=11,color=GREY,italic=True)
notes(s,"This is the Marts layer (Phase 2, downstream of Core via dbt). Pitch = highest-value case of the OMOP 'write analytics once' thesis: derived variables defined once, canonically, with lineage — vs. re-coded per study (the 3 pain repos). dbt features map to our principles: dbt docs DAG = lineage to source (#12); SQL is the definition (no black box, #12); source()/ref()+access enforce Core→Analytic→Marts (#8); tests move much of the 7,025-rule QC into CI; sensitivity_tier flows via meta, enforced in IAM (#11); model versions/exposures = a derived-variable catalog. Recommend one mart PER construct (grouped), not a wide mega-table. Grain: participant, or participant×wave for longitudinal. Examples map to established cancer-prevention exposure domains (WCRF/AICR). See CLAUDE.md — Derived variables, lineage & dbt / Mart catalog.")

# 10 QUERIES
s=slide();titlebar(s,"Same queries, three ways","From painful/impossible → routine → trivial")
table(s,[["Standard query","Wide (today)","Phase 1","Phase 2"],
 ["Multi-select distribution","unpivot + v1/v2 COALESCE","filtered group-by; versions pool","plain group-by / view"],
 ["Labeled distribution, any question","bespoke CASE, no reuse","parameterized by concept_id","precomputed aggregate"],
 ["Completion & true missingness","ambiguous","still ambiguous","sessions + skip logic"],
 ["Non-PHI extract (governance)","manual allow-list","manual allow-list","sensitivity tier + IAM"],
 ["Harmonized field across surveys","rebuild crosswalk (3 repos)","pool a reused concept","equivalence plane (planned)"]],
 Inches(0.4),Inches(1.5),Inches(12.5),Inches(3.6),col_widths=[Inches(3.5),Inches(3.0),Inches(3.0),Inches(3.0)],fs=12.5,hdr_fs=13)
textbox(s,"Phase 1 fixes the loud, everyday pains; Phase 2 unlocks missingness, governance & harmonization.",Inches(0.4),Inches(6.4),Inches(12.5),Inches(0.6),size=15,color=GREY,italic=True)
notes(s,"Worked SQL for each row is in internal_pitch.md — Value proposition (Q1–Q5).")

# 11 RECOMMENDATION
s=slide();titlebar(s,"Recommendation")
bullets(s,["Approve Phase 1 now — low cost, low risk, immediately useful; its responses fact carries into Phase 2 unchanged",
 "Commit to Phase 2 — governance + lineage; non-negotiable for sharing data externally through PR2"],
 Inches(0.8),Inches(1.8),Inches(11.7),Inches(2.4),size=21,gap=18)
box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(4.6),Inches(11.7),Inches(1.4))
box.fill.solid();box.fill.fore_color.rgb=LIGHT;box.line.color.rgb=TEAL;box.shadow.inherit=False
tf=box.text_frame;tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE;tf.margin_left=Inches(0.3)
p=tf.paragraphs[0];r=p.add_run();r.text="Phase 1 = an internal quick win that proves the model.   Phase 2 = the governed, shareable product."
r.font.size=Pt(18);r.font.bold=True;r.font.color.rgb=NAVY;r.font.name="Calibri"
notes(s,"Ask: approve Phase 1 now; commit Phase 2 (esp. governance); consider pulling the concept-equivalence plane forward (geocoding). See internal_pitch.md — Recommendation.")

out="docs/Connect_Data_Model_Pitch.pptx";prs.save(out)
print("saved",out,"with",len(prs.slides._sldIdLst),"slides")
